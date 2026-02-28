#!/usr/bin/env python3
"""
Create PowerPoint deck combining all market analysis slides
US Agency Market Analysis - Complete Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

print("="*80)
print("CREATING POWERPOINT PRESENTATION")
print("="*80)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# ============================================================================
# SLIDE 1: TITLE SLIDE
# ============================================================================
print("\nAdding Slide 1: Title Slide")

slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(15), Inches(2))
title_frame = title_box.text_frame
title_frame.text = "US Agency Market Analysis"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(72)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(26, 26, 26)
title_para.alignment = PP_ALIGN.CENTER

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(15), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "2023 Market Overview by Region, Port & Vessel Class"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(32)
subtitle_para.font.color.rgb = RGBColor(102, 102, 102)
subtitle_para.alignment = PP_ALIGN.CENTER

# Add market stats
stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(15), Inches(1))
stats_frame = stats_box.text_frame
stats_frame.text = "Total Market: $167.7M | 49,726 Port Calls"
stats_para = stats_frame.paragraphs[0]
stats_para.font.size = Pt(24)
stats_para.font.color.rgb = RGBColor(52, 152, 219)
stats_para.alignment = PP_ALIGN.CENTER

# Add background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# ============================================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================================
print("Adding Slide 2: Executive Summary")

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "Executive Summary"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(26, 26, 26)

# Content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(6.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

summary_text = """Market Overview
• Total US Agency Market: $167.7M in annual revenue
• 49,726 port calls across all US ports
• Average fee per call: $3,373

Regional Concentration
• Gulf Coast dominates: 47.8% of total market revenue
• Houston is the mega-market: $24.7M (17.6% of US total)
• Top 5 ports = 46.9% of market | Top 10 ports = 72.0% of market

Vessel Class Insights
• Bulk Carriers: 38.2% of revenue ($64.0M) with $9,000 avg fee
• Tankers (Oil/Crude): 25.4% of revenue ($42.5M) with $4,500 avg fee
• Container Ships: High volume (28% of calls) but low revenue (5.4%) at $650/call

Strategic Implications
• Market highly concentrated in Gulf Coast energy/bulk trade
• Texas triangle (Houston + South Texas) = 31.4% of entire US market
• Container shipping represents volume opportunity with lower margins
"""

content_frame.text = summary_text
for paragraph in content_frame.paragraphs:
    paragraph.font.size = Pt(16)
    paragraph.font.color.rgb = RGBColor(51, 51, 51)
    paragraph.space_before = Pt(6)

# ============================================================================
# SLIDE 3: REGIONAL MARKET ANALYSIS (First map slide)
# ============================================================================
print("Adding Slide 3: Regional Market Analysis")

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

img_path = r'G:\My Drive\LLM\task_usace_entrance_clearance\00_DATA\00.03_REPORTS\SLIDE_Regional_Market_Analysis.png'
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16))
    print(f"  Added: {os.path.basename(img_path)}")

# ============================================================================
# SLIDE 4: GEOGRAPHIC MARKET MAP (Full US map)
# ============================================================================
print("Adding Slide 4: Geographic Market Map")

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

img_path = r'G:\My Drive\LLM\task_usace_entrance_clearance\00_DATA\00.03_REPORTS\SLIDE_Geographic_Market_Map.png'
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16))
    print(f"  Added: {os.path.basename(img_path)}")

# ============================================================================
# SLIDE 5: REGIONAL DEEP DIVE (4 zoom maps)
# ============================================================================
print("Adding Slide 5: Regional Deep Dive")

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

img_path = r'G:\My Drive\LLM\task_usace_entrance_clearance\00_DATA\00.03_REPORTS\SLIDE_Regional_Geographic_Zooms.png'
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16))
    print(f"  Added: {os.path.basename(img_path)}")

# ============================================================================
# SLIDE 6: REVENUE DENSITY HEATMAP
# ============================================================================
print("Adding Slide 6: Revenue Density Heatmap")

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

img_path = r'G:\My Drive\LLM\task_usace_entrance_clearance\00_DATA\00.03_REPORTS\SLIDE_Revenue_Density_Heatmap.png'
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16))
    print(f"  Added: {os.path.basename(img_path)}")

# ============================================================================
# SLIDE 7: TOP PORTS ANALYSIS
# ============================================================================
print("Adding Slide 7: Top Ports Analysis")

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

img_path = r'G:\My Drive\LLM\task_usace_entrance_clearance\00_DATA\00.03_REPORTS\SLIDE_Top_Ports_Analysis.png'
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16))
    print(f"  Added: {os.path.basename(img_path)}")

# ============================================================================
# SLIDE 8: KEY TAKEAWAYS
# ============================================================================
print("Adding Slide 8: Key Takeaways")

slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "Key Takeaways & Strategic Insights"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(26, 26, 26)

# Create two columns
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(6.5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

left_text = """Geographic Concentration
• Gulf Coast = 48% of US market
• Houston mega-market dominates
• Texas corridor (Houston to Brownsville)
  controls 31% of total revenue

Market Structure
• Highly concentrated: Top 10 ports = 72%
• Energy & bulk commodities drive revenue
• Container = high volume, low margin
• Tankers/bulk = lower volume, high margin

Regional Champions
• North Texas: Houston ($24.7M)
• South Atlantic: South Florida ($12.6M)
• Gulf East: New Orleans ($8.6M)
• Mid-Atlantic: New York ($9.2M)
• California: LA-Long Beach ($7.3M)
"""

left_frame.text = left_text
for paragraph in left_frame.paragraphs:
    paragraph.font.size = Pt(14)
    paragraph.font.color.rgb = RGBColor(51, 51, 51)
    paragraph.space_before = Pt(4)

right_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.5), Inches(7), Inches(6.5))
right_frame = right_box.text_frame
right_frame.word_wrap = True

right_text = """Vessel Class Opportunities
• Bulk Carriers: Highest revenue/call
  ($9,000 avg) - 38% of market
• Tankers: Strong margin ($4,500 avg)
  Oil/Chemical = 39% combined
• Gas Carriers: Premium segment
  LNG/LPG = 8.2% of market
• Container: Volume play ($650 avg)
  but 28% of all port calls

Market Entry Strategy
• Focus on Gulf Coast for bulk/energy
• East Coast for container volume
• Regional approach vs. national
• Partner with established ports
• Specialize by vessel class

Growth Potential
• Container volume growing
• Energy export expansion (LNG)
• Chemical tanker demand increasing
• Regional consolidation trends
"""

right_frame.text = right_text
for paragraph in right_frame.paragraphs:
    paragraph.font.size = Pt(14)
    paragraph.font.color.rgb = RGBColor(51, 51, 51)
    paragraph.space_before = Pt(4)

# ============================================================================
# SAVE PRESENTATION
# ============================================================================

output_file = r'G:\My Drive\LLM\task_usace_entrance_clearance\00_DATA\00.03_REPORTS\US_Agency_Market_Analysis_2023.pptx'
prs.save(output_file)

print("\n" + "="*80)
print("POWERPOINT DECK CREATED SUCCESSFULLY")
print("="*80)
print(f"File: {output_file}")
print(f"Total Slides: {len(prs.slides)}")
print("="*80)
print("\nSlide Order:")
print("  1. Title Slide")
print("  2. Executive Summary")
print("  3. Regional Market Analysis (Bubble Map)")
print("  4. Geographic Market Map (Full US)")
print("  5. Regional Deep Dive (4 Zooms)")
print("  6. Revenue Density Heatmap")
print("  7. Top Ports Analysis")
print("  8. Key Takeaways & Strategic Insights")
print("="*80)
