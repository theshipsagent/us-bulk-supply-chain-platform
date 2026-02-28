"""
Build a professional PowerPoint deck for the Lower Mississippi River
Industrial Infrastructure Report — styled to match oceandatum.ai.

Dark theme:  #0a0a0a background, #64ffb4 accent, white text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = r"G:\My Drive\LLM\sources_data_maps\_html_web_files\Lower_Mississippi_Infrastructure_Deck_v2.pptx"

# ── Brand colors (oceandatum.ai) ────────────────────────────────────
BG       = RGBColor(0x0a, 0x0a, 0x0a)
BG_CARD  = RGBColor(0x14, 0x14, 0x18)
ACCENT   = RGBColor(0x64, 0xff, 0xb4)
ACCENT2  = RGBColor(0x0e, 0xa5, 0xe9)   # secondary blue
WHITE    = RGBColor(0xff, 0xff, 0xff)
GRAY     = RGBColor(0x99, 0x99, 0x99)
DIM      = RGBColor(0x66, 0x66, 0x66)
RED_ACC  = RGBColor(0xe9, 0x45, 0x60)
ORANGE   = RGBColor(0xff, 0x6b, 0x35)
CYAN     = RGBColor(0x0e, 0xa5, 0xe9)
GREEN    = RGBColor(0x22, 0xc5, 0x5e)
PURPLE   = RGBColor(0xa8, 0x55, 0xf7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank

# ── Helpers ──────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.15, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    try:
        p.line_spacing = Pt(font_size * line_spacing)
    except Exception:
        pass
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT, height=Pt(3)):
    line = add_rect(slide, left, top, width, height, color)
    return line

def add_stat_card(slide, left, top, number, label, accent_color=ACCENT):
    """Add a big-number stat card."""
    card_w = Inches(2.8)
    card_h = Inches(2.0)
    card = add_rect(slide, left, top, card_w, card_h, BG_CARD,
                    border_color=RGBColor(0x2a, 0x2a, 0x3a), border_width=Pt(1))
    # Accent bar at top of card
    add_rect(slide, left, top, card_w, Pt(3), accent_color)
    # Big number
    add_text_box(slide, left + Inches(0.25), top + Inches(0.25),
                 card_w - Inches(0.5), Inches(1.0),
                 number, font_size=42, color=accent_color, bold=True,
                 alignment=PP_ALIGN.LEFT)
    # Label
    add_text_box(slide, left + Inches(0.25), top + Inches(1.2),
                 card_w - Inches(0.5), Inches(0.7),
                 label, font_size=13, color=GRAY, bold=False,
                 alignment=PP_ALIGN.LEFT)

def add_branded_footer(slide):
    """Add oceandatum.ai branding footer."""
    # Thin line
    add_rect(slide, Inches(0.6), SLIDE_H - Inches(0.55),
             SLIDE_W - Inches(1.2), Pt(0.5),
             RGBColor(0x2a, 0x2a, 0x3a))
    # Brand text
    add_text_box(slide, Inches(0.6), SLIDE_H - Inches(0.50),
                 Inches(4), Inches(0.35),
                 "oceandatum.ai  |  DATUM Intelligence",
                 font_size=9, color=DIM, bold=False)
    # Date
    add_text_box(slide, SLIDE_W - Inches(3.6), SLIDE_H - Inches(0.50),
                 Inches(3), Inches(0.35),
                 "February 2026  |  Confidential",
                 font_size=9, color=DIM, bold=False,
                 alignment=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1)

# Large accent line at top
add_rect(slide1, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT)

# Title
add_text_box(slide1, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
             "Lower Mississippi River", font_size=48, color=WHITE, bold=True,
             font_name="Calibri Light")

add_text_box(slide1, Inches(0.8), Inches(2.75), Inches(11), Inches(0.9),
             "Industrial Infrastructure Report", font_size=40, color=ACCENT, bold=False,
             font_name="Calibri Light")

# Accent bar under title
add_accent_line(slide1, Inches(0.8), Inches(3.65), Inches(3.5), ACCENT, Pt(3))

# Subtitle
add_text_box(slide1, Inches(0.8), Inches(4.0), Inches(9), Inches(0.5),
             "Comprehensive Survey of Facilities from Head of Passes to Baton Rouge  |  Mile 0 to Mile 232",
             font_size=16, color=GRAY, bold=False)

# Key stat line
add_text_box(slide1, Inches(0.8), Inches(4.7), Inches(9), Inches(0.5),
             "112 Facilities  \u2022  7 Refineries  \u2022  2M+ BPD Refining Capacity  \u2022  43 Pipelines  \u2022  5 Port Authorities",
             font_size=13, color=DIM, bold=False)

# Bottom right branding
add_text_box(slide1, Inches(9.5), Inches(5.8), Inches(3.5), Inches(0.4),
             "oceandatum.ai", font_size=22, color=ACCENT, bold=False,
             alignment=PP_ALIGN.RIGHT, font_name="Georgia")
add_text_box(slide1, Inches(9.5), Inches(6.2), Inches(3.5), Inches(0.3),
             "Charting Clarity from Complexity", font_size=11, color=DIM, bold=False,
             alignment=PP_ALIGN.RIGHT, font_name="Georgia")
add_text_box(slide1, Inches(9.5), Inches(6.55), Inches(3.5), Inches(0.3),
             "February 2026", font_size=11, color=DIM, bold=False,
             alignment=PP_ALIGN.RIGHT)

add_branded_footer(slide1)


# =====================================================================
# SLIDE 2: CORRIDOR AT A GLANCE
# =====================================================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2)

add_rect(slide2, Inches(0), Inches(0), SLIDE_W, Pt(4), ACCENT)

# Section title
add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "CORRIDOR AT A GLANCE", font_size=11, color=ACCENT, bold=True)
add_text_box(slide2, Inches(0.8), Inches(0.75), Inches(10), Inches(0.5),
             "232-Mile Industrial Corridor  \u2014  Head of Passes to Baton Rouge",
             font_size=26, color=WHITE, bold=True, font_name="Calibri Light")

add_accent_line(slide2, Inches(0.8), Inches(1.35), Inches(2.5))

# Row 1 of stat cards
y1 = Inches(1.7)
gap = Inches(3.1)
add_stat_card(slide2, Inches(0.8),          y1, "2.02M", "Barrels/Day\nRefining Capacity (6 active)", ACCENT)
add_stat_card(slide2, Inches(0.8) + gap,    y1, "112",   "Industrial Facilities\nConsolidated from 180 zones", ACCENT2)
add_stat_card(slide2, Inches(0.8) + gap*2,  y1, "43",    "Pipeline Systems\n24 Crude  |  11 HGL  |  8 Product", PURPLE)
add_stat_card(slide2, Inches(0.8) + gap*3,  y1, "5",     "Port Authorities\nGoverning Jurisdiction", GREEN)

# Row 2 of stat cards
y2 = Inches(4.1)
add_stat_card(slide2, Inches(0.8),          y2, "7",     "Petroleum Refineries\nIncl. 4th & 6th largest in U.S.", RED_ACC)
add_stat_card(slide2, Inches(0.8) + gap,    y2, "16",    "Chemical Plants\nPetrochemical & fertilizer", ORANGE)
add_stat_card(slide2, Inches(0.8) + gap*2,  y2, "32",    "Tank Storage Terminals\nCrude, refined products, chemicals", CYAN)
add_stat_card(slide2, Inches(0.8) + gap*3,  y2, "9",     "Export Grain Elevators\nGlobal ag commodity gateway", GREEN)

add_branded_footer(slide2)


# =====================================================================
# SLIDE 3: REFINERY CORRIDOR
# =====================================================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3)

add_rect(slide3, Inches(0), Inches(0), SLIDE_W, Pt(4), RED_ACC)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "REFINERY CORRIDOR", font_size=11, color=RED_ACC, bold=True)
add_text_box(slide3, Inches(0.8), Inches(0.75), Inches(10), Inches(0.5),
             "Seven Refineries  \u2014  Combined 2.02 Million BPD Capacity",
             font_size=26, color=WHITE, bold=True, font_name="Calibri Light")
add_accent_line(slide3, Inches(0.8), Inches(1.35), Inches(2.5), RED_ACC)

# Table of refineries
refineries = [
    ("Valero Meraux",           "Mile 87",  "Valero Energy",    "135,000",  "Medium-heavy sour",      "Active"),
    ("PBF Chalmette",           "Mile 89",  "PBF Energy",       "185,000",  "Light sweet to heavy sour","Active"),
    ("Valero St. Charles",      "Mile 125", "Valero Energy",    "340,000",  "Sour crudes via LOOP",    "Active"),
    ("Shell Norco",             "Mile 126", "Shell USA",        "250,000",  "Medium sweet crude",      "Active"),
    ("Marathon Garyville",      "Mile 140", "Marathon Petroleum","596,000", "Heavy sour (4th largest US)","Active"),
    ("Shell Convent",           "Mile 168", "Shell USA",        "211,000",  "N/A \u2014 closed 2020",  "Converting"),
    ("ExxonMobil Baton Rouge",  "Mile 232", "ExxonMobil",       "522,500", "Broad flexibility",       "Active"),
]

# Table header
tbl_left = Inches(0.8)
tbl_top = Inches(1.65)
col_widths = [Inches(2.6), Inches(1.0), Inches(2.2), Inches(1.5), Inches(2.8), Inches(1.1)]
headers = ["Facility", "Mile", "Owner", "Capacity (BPD)", "Crude Slate", "Status"]
row_h = Inches(0.38)

# Header row bg
total_w = sum(c for c in col_widths)
add_rect(slide3, tbl_left, tbl_top, total_w, row_h, RGBColor(0x1a, 0x1a, 0x2e))

x = tbl_left
for i, hdr in enumerate(headers):
    add_text_box(slide3, x, tbl_top, col_widths[i], row_h,
                 hdr, font_size=10, color=ACCENT, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    x += col_widths[i]

# Data rows
for r, ref in enumerate(refineries):
    y = tbl_top + row_h * (r + 1)
    row_bg = BG_CARD if r % 2 == 0 else BG
    add_rect(slide3, tbl_left, y, total_w, row_h, row_bg)
    # Subtle bottom border
    add_rect(slide3, tbl_left, y + row_h - Pt(0.5), total_w, Pt(0.5),
             RGBColor(0x2a, 0x2a, 0x3a))
    x = tbl_left
    for c, val in enumerate(ref):
        clr = WHITE
        bld = False
        if c == 0:
            bld = True
        elif c == 3:
            clr = ACCENT
            bld = True
        elif c == 5:
            clr = GREEN if val == "Active" else ORANGE
        add_text_box(slide3, x, y, col_widths[c], row_h,
                     val, font_size=10.5, color=clr, bold=bld,
                     anchor=MSO_ANCHOR.MIDDLE)
        x += col_widths[c]

# Summary bar at bottom of table
summary_y = tbl_top + row_h * 8
add_rect(slide3, tbl_left, summary_y, total_w, Inches(0.45),
         RGBColor(0x14, 0x14, 0x1e),
         border_color=RED_ACC, border_width=Pt(1))
add_text_box(slide3, tbl_left + Inches(0.3), summary_y, total_w - Inches(0.6), Inches(0.45),
             "Combined Operating Capacity (6 Active):  2,018,500 BPD   \u2022   "
             "4th largest (Marathon Garyville) & 6th largest (ExxonMobil BR) refineries in the U.S.",
             font_size=11, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

# Key connections callout
callout_y = Inches(5.6)
add_text_box(slide3, Inches(0.8), callout_y, Inches(3), Inches(0.3),
             "KEY INFRASTRUCTURE CONNECTIONS", font_size=9, color=DIM, bold=True)
add_accent_line(slide3, Inches(0.8), callout_y + Inches(0.3), Inches(1.5), DIM, Pt(1))

connections = [
    "LOOP (LA Offshore Oil Port) \u2014 deepwater crude imports for 5 of 6 active refineries",
    "Capline Pipeline \u2014 reversed 2021, now delivers 417K bpd southbound from Patoka, IL to St. James hub",
    "Colonial & Plantation Pipelines \u2014 product distribution to Southeast & East Coast markets",
]
for i, conn in enumerate(connections):
    add_text_box(slide3, Inches(1.0), callout_y + Inches(0.4) + Inches(i * 0.32),
                 Inches(11), Inches(0.3),
                 "\u25B8  " + conn, font_size=10.5, color=GRAY)

add_branded_footer(slide3)


# =====================================================================
# SLIDE 4: INFRASTRUCTURE BREAKDOWN
# =====================================================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4)

add_rect(slide4, Inches(0), Inches(0), SLIDE_W, Pt(4), CYAN)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "INFRASTRUCTURE BREAKDOWN", font_size=11, color=CYAN, bold=True)
add_text_box(slide4, Inches(0.8), Inches(0.75), Inches(10), Inches(0.5),
             "112 Facilities Across 8 Infrastructure Categories",
             font_size=26, color=WHITE, bold=True, font_name="Calibri Light")
add_accent_line(slide4, Inches(0.8), Inches(1.35), Inches(2.5), CYAN)

# Category cards — 2 columns x 4 rows
categories = [
    ("Refineries",       "7",  RED_ACC,  "2.02M bpd combined capacity\n6 active, 1 converting to alt fuels\nMarathon Garyville = 4th largest in US"),
    ("Chemical Plants",  "16", ORANGE,   "Petrochemicals, fertilizers, chlor-alkali\nCF Industries = largest N\u2082 complex in world\nDow, Shell, OxyChem, BASF, Shintech"),
    ("Tank Storage",     "32", CYAN,     "Crude, products, chemicals, LPG\nIMTT, Kinder Morgan, MPLX, NuStar\nSt. James hub = physical crude pricing point"),
    ("Grain Elevators",  "9",  GREEN,    "ADM, Cargill, Bunge, Zen-Noh, CHS\n60%+ of US grain exports transit here\nMGMT = closest elevator to Gulf"),
    ("Bulk Terminals",   "4",  PURPLE,   "IMT Myrtle Grove, United Bulk, ABT\nCoal, petcoke, iron ore, limestone\nDeepwater ship-loading capability"),
    ("Bulk Plants",      "9",  RGBColor(0x8b,0x5c,0xf6), "Alumina, steel, sugar, biomass\nNoranda Alumina, Nucor Steel, Domino\nSpecialized processing + river access"),
    ("General Cargo",    "15", RGBColor(0xea,0xb3,0x08), "Port NOLA wharves, Napoleon Ave containers\nNashville Ave = busiest by vessel events\nJulia & Erato = cruise terminals"),
    ("Mid-Stream",       "20", GRAY,     "Ship-to-barge transfer buoy fields\nARTCO, Cooper/Darrow, Associated\n13 operators across Mile 57\u2013180"),
]

card_w = Inches(5.8)
card_h = Inches(1.25)
col1_x = Inches(0.8)
col2_x = Inches(7.0)
start_y = Inches(1.65)
row_gap = Inches(1.35)

for i, (cat_name, count, color, desc) in enumerate(categories):
    col = i % 2
    row = i // 2
    x = col1_x if col == 0 else col2_x
    y = start_y + row * row_gap

    # Card background
    add_rect(slide4, x, y, card_w, card_h, BG_CARD,
             border_color=RGBColor(0x2a, 0x2a, 0x3a), border_width=Pt(1))
    # Left accent bar
    add_rect(slide4, x, y, Pt(4), card_h, color)

    # Count
    add_text_box(slide4, x + Inches(0.2), y + Inches(0.05),
                 Inches(0.7), card_h - Inches(0.1),
                 count, font_size=32, color=color, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    # Category name
    add_text_box(slide4, x + Inches(0.85), y + Inches(0.08),
                 Inches(2.0), Inches(0.3),
                 cat_name, font_size=13, color=WHITE, bold=True)
    # Description
    add_text_box(slide4, x + Inches(0.85), y + Inches(0.38),
                 card_w - Inches(1.1), Inches(0.85),
                 desc, font_size=9.5, color=GRAY, line_spacing=1.4)

add_branded_footer(slide4)


# =====================================================================
# SLIDE 5: TRANSPORT NETWORK & CONNECTIVITY
# =====================================================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5)

add_rect(slide5, Inches(0), Inches(0), SLIDE_W, Pt(4), PURPLE)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "TRANSPORT NETWORK", font_size=11, color=PURPLE, bold=True)
add_text_box(slide5, Inches(0.8), Inches(0.75), Inches(10), Inches(0.5),
             "Pipeline, Rail & Maritime Connectivity",
             font_size=26, color=WHITE, bold=True, font_name="Calibri Light")
add_accent_line(slide5, Inches(0.8), Inches(1.35), Inches(2.5), PURPLE)

# Three columns: Pipelines | Rail | Maritime
col_w = Inches(3.7)
col_gap = Inches(0.35)
col_h = Inches(4.8)
col_y = Inches(1.7)

for i, (title, accent, items) in enumerate([
    ("PIPELINES (43 Systems)", RGBColor(0xff, 0x44, 0x44), [
        ("Crude Oil", "24 systems", "#ff4444",
         "Capline (417K bpd, reversed 2021)\n"
         "LOOP subsea pipeline to Clovelly\n"
         "Zydeco (4 segments), Mars, Bayou Bridge\n"
         "Delta, Genesis, ETCOP systems"),
        ("HGL / NGL", "11 systems", "#ff44aa",
         "ATEX Express, Aegis Ethane\n"
         "Cajun-Sibon, Lou-Tex, TEPPCO\n"
         "Bayou Ethane, Dixie, Tri-States"),
        ("Refined Products", "8 systems", "#44aaff",
         "Colonial (Lines 1 & 2), Plantation\n"
         "Bengal, Parkway (Valero), Explorer\n"
         "Centennial, Collins"),
    ]),
    ("RAIL (8 Carriers)", RGBColor(0xff, 0xe1, 0x19), [
        ("Class I Railroads", "6 carriers", "#ffe119",
         "Union Pacific, BNSF Railway\n"
         "Norfolk Southern, CSX Transportation\n"
         "Canadian National, CPKC"),
        ("Terminal Railroads", "2 carriers", "#e6194b",
         "New Orleans Public Belt (NOPB)\n"
         "New Orleans & Gulf Coast (NOGC)\n"
         "Interchange + last-mile service"),
        ("Key Routes", "", "#aaaaaa",
         "UP & BNSF main lines parallel river\n"
         "NS/CSX interchange at New Orleans\n"
         "On-dock rail at Napoleon Ave terminal"),
    ]),
    ("MARITIME", RGBColor(0x06, 0xb6, 0xd4), [
        ("Deepwater Access", "", "#06b6d4",
         "LOOP: supertanker crude imports\n"
         "SW Pass: 45-ft draft channel\n"
         "Port Fourchon: GoM service base"),
        ("Port Authorities", "5 jurisdictions", "#64ffb4",
         "Port of Greater Baton Rouge (RM 169\u2013253)\n"
         "Port of South Louisiana (RM 115\u2013169)\n"
         "Port of New Orleans (RM 82\u2013115)\n"
         "St. Bernard Port | Plaquemines Port"),
        ("Vessel Activity", "", "#aaaaaa",
         "Nashville Ave: 5,955+ vessel events\n"
         "Napoleon Ave: 366K TEU capacity\n"
         "20 midstream buoy transfer points"),
    ]),
]):
    cx = Inches(0.8) + i * (col_w + col_gap)

    # Column card
    add_rect(slide5, cx, col_y, col_w, col_h, BG_CARD,
             border_color=RGBColor(0x2a, 0x2a, 0x3a), border_width=Pt(1))
    # Top accent bar
    add_rect(slide5, cx, col_y, col_w, Pt(3), accent)

    # Column title
    add_text_box(slide5, cx + Inches(0.2), col_y + Inches(0.15),
                 col_w - Inches(0.4), Inches(0.3),
                 title, font_size=12, color=accent, bold=True)

    # Sub-sections
    for j, (sub_title, sub_count, sub_clr, sub_desc) in enumerate(items):
        sy = col_y + Inches(0.55) + j * Inches(1.45)
        # Sub-header
        label = sub_title
        if sub_count:
            label += f"  ({sub_count})"
        hx = sub_clr.lstrip('#')
        sc = RGBColor(int(hx[0:2],16), int(hx[2:4],16), int(hx[4:6],16))
        add_text_box(slide5, cx + Inches(0.2), sy,
                     col_w - Inches(0.4), Inches(0.25),
                     label, font_size=10, color=sc, bold=True)
        # Thin divider
        add_rect(slide5, cx + Inches(0.2), sy + Inches(0.25),
                 Inches(1.2), Pt(0.5), RGBColor(0x33, 0x33, 0x44))
        # Description
        add_text_box(slide5, cx + Inches(0.2), sy + Inches(0.32),
                     col_w - Inches(0.4), Inches(1.0),
                     sub_desc, font_size=9.5, color=GRAY, line_spacing=1.35)

add_branded_footer(slide5)


# ── Speaker Notes ────────────────────────────────────────────────────
slide1.notes_slide.notes_text_frame.text = (
    "Lower Mississippi River Industrial Infrastructure Report — February 2026.\n"
    "This deck summarizes a comprehensive facility-by-facility survey of the 232-mile industrial corridor "
    "from Head of Passes (Mile 0, Southwest Pass) to Baton Rouge (Mile 232).\n"
    "Data sourced from USDA AMS/FGIS, USACE, company filings, EIA, port authority documents, "
    "and trade publications. 112 facilities consolidated from 180 individual berths/zones.\n"
    "Interactive map and full report available at oceandatum.ai."
)

slide2.notes_slide.notes_text_frame.text = (
    "Corridor At A Glance — Key Statistics:\n"
    "- 2.02 million BPD combined refining capacity across 6 active refineries (7th facility, Shell Convent, "
    "closed as a refinery in Dec 2020 and is converting to alternative fuels).\n"
    "- 112 industrial facilities spanning 8 categories: refineries, chemical plants, tank storage, "
    "grain elevators, bulk terminals, bulk plants, general cargo/wharves, and midstream transfer points.\n"
    "- 43 pipeline systems connecting the corridor: 24 crude oil, 11 HGL/NGL, 8 refined product.\n"
    "- 5 port authorities govern jurisdiction from Head of Passes to north of Baton Rouge.\n"
    "- This corridor handles approximately 60% of U.S. grain exports and is the largest concentration "
    "of petrochemical and refining infrastructure on the Gulf Coast."
)

slide3.notes_slide.notes_text_frame.text = (
    "Refinery Corridor Detail:\n"
    "- Marathon Garyville (596K bpd) is the 4th largest refinery in the U.S. Heavy sour crude specialist. "
    "$3.9B Garyville Major Expansion completed 2009.\n"
    "- ExxonMobil Baton Rouge (522K bpd) is the 6th largest in the U.S. and 17th worldwide. "
    "Integrated with the adjacent chemical complex — 300+ products and grades.\n"
    "- Valero operates two facilities (Meraux 135K + St. Charles 340K = 475K bpd combined). "
    "St. Charles is investing $230M in FCC upgrade for completion H2 2026.\n"
    "- Shell Norco (250K bpd) is Shell's sole U.S. Energy & Chemicals Park — integrated refinery + chemical plant.\n"
    "- PBF Chalmette (185K bpd) has Nelson Complexity Index of 13.0, dual-coker configuration.\n"
    "- Shell Convent (211K bpd former capacity) permanently closed Dec 2020. Converting to renewable diesel "
    "via the Peri Project — $1.5B in approved Louisiana tax relief.\n"
    "- LOOP provides deepwater crude access for 5 of 6 active refineries. Capline reversal (2021) now delivers "
    "417K bpd southbound from Patoka, IL to the St. James storage hub."
)

slide4.notes_slide.notes_text_frame.text = (
    "Infrastructure Breakdown by Category:\n"
    "- Chemical Plants (16): CF Industries Donaldsonville is the world's largest nitrogen fertilizer complex "
    "(~8M tons/year). Cos-Mar/TotalEnergies Carville is the world's largest styrene plant (~1.2M tons/year). "
    "OxyChem Taft acquired by Berkshire Hathaway Jan 2026 ($9.7B).\n"
    "- Tank Storage (32): St. James storage hub is the physical delivery point for NYMEX Light Sweet Crude. "
    "Major operators include IMTT (4 terminals), Kinder Morgan (3), MPLX (2), NuStar, Plains, Capline.\n"
    "- Grain Elevators (9): ADM (3 locations), Cargill (2), Bunge, Zen-Noh, CHS, Louis Dreyfus. "
    "MGMT at Myrtle Grove (Mile 57) is the closest export elevator to the Gulf.\n"
    "- General Cargo (15): Napoleon Ave container terminal expanding from 366K to 1.5M TEU capacity. "
    "Louisiana International Terminal ($1.8B) planned at Violet Dock — construction 2025, first berth 2028."
)

slide5.notes_slide.notes_text_frame.text = (
    "Transport Network & Connectivity:\n"
    "- Pipelines: 43 systems mapped. Capline Pipeline (Marathon/MPLX) reversed in 2021, now the primary "
    "crude artery delivering 417K bpd southbound. Colonial and Plantation pipelines distribute refined "
    "products to the Southeast and East Coast. LOOP subsea pipeline connects to Clovelly storage.\n"
    "- Rail: 6 Class I carriers (UP, BNSF, NS, CSXT, CN, CPKC) plus 2 terminal railroads (NOPB, NOGC). "
    "NOPB provides interchange and last-mile switching for Port NOLA facilities. "
    "Napoleon Ave terminal has on-dock intermodal rail ramp.\n"
    "- Maritime: Southwest Pass provides 45-ft draft channel access. LOOP handles supertanker crude imports "
    "18 miles offshore. Port Fourchon supports 90%+ of deepwater Gulf of Mexico operations. "
    "5 port authorities: Greater Baton Rouge, South Louisiana, New Orleans, St. Bernard, Plaquemines.\n"
    "- Nashville Ave wharves recorded 5,955+ vessel events — busiest facility by traffic on the corridor."
)


# ── Save ─────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
prs.save(OUTPUT)
size_kb = os.path.getsize(OUTPUT) / 1024
print(f"Deck written: {OUTPUT}")
print(f"Size: {size_kb:.0f} KB")
print(f"Slides: {len(prs.slides)}")
