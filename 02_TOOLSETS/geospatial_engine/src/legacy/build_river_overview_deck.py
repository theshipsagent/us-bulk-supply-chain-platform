"""
Mississippi River Facilities Overview Presentation
Creates comprehensive slide deck with facility types, vessels, cargoes, and maps
"""

import pandas as pd
import folium
from folium.plugins import MarkerCluster
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from PIL import Image
import base64

# Read facility data
facilities = pd.read_csv('master_facilities.csv')

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(18, 48, 87)  # Navy blue

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = txBox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    tf = txBox.text_frame
    tf.text = subtitle
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 220, 255)

    return slide

def add_content_slide(prs, title, layout_num=5):
    """Add content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_num])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(18, 48, 87)
    return slide

def create_facility_map(facilities_df, zoom_start=7):
    """Create folium map of facilities"""
    # Center on Louisiana
    center_lat = facilities_df['lat'].mean()
    center_lon = facilities_df['lng'].mean()

    m = folium.Map(
        location=[center_lat, center_lon],
        zoom_start=zoom_start,
        tiles='CartoDB positron'
    )

    # Color mapping for facility types
    colors = {
        'Tank Storage': '#e74c3c',  # Red
        'Refinery': '#c0392b',  # Dark red
        'Bulk Terminal': '#3498db',  # Blue
        'Elevator': '#f39c12',  # Orange
        'General Cargo': '#27ae60',  # Green
        'Mid-Stream': '#9b59b6',  # Purple
        'Bulk Plant': '#16a085',  # Teal
    }

    # Add facilities
    marker_cluster = MarkerCluster().add_to(m)

    for idx, row in facilities_df.iterrows():
        if pd.notna(row['lat']) and pd.notna(row['lng']):
            color = colors.get(row['facility_type'], '#95a5a6')

            popup_html = f"""
            <div style="width: 300px;">
                <h4>{row['canonical_name']}</h4>
                <b>Type:</b> {row['facility_type']}<br>
                <b>Location:</b> {row['city']}, {row['state']}<br>
                <b>River Mile:</b> {row['mrtis_mile']}<br>
                <b>Commodities:</b> {row['commodities'][:100] if pd.notna(row['commodities']) else 'N/A'}...
            </div>
            """

            folium.CircleMarker(
                location=[row['lat'], row['lng']],
                radius=6,
                popup=folium.Popup(popup_html, max_width=300),
                color=color,
                fill=True,
                fillColor=color,
                fillOpacity=0.7,
                weight=2
            ).add_to(marker_cluster)

    return m

def save_map_as_image(folium_map, filename, width=1200, height=800):
    """Save folium map as PNG image"""
    # Save map as HTML
    html_file = f"{filename}.html"
    folium_map.save(html_file)

    # For now, return the HTML file path
    # In production, you'd use selenium or similar to convert HTML to PNG
    return html_file

# SLIDE 1: Title
add_title_slide(
    prs,
    "Mississippi River Infrastructure",
    "Comprehensive Overview of Facilities, Vessels & Cargo Operations"
)

# SLIDE 2: System Overview
slide = add_content_slide(prs, "Mississippi River System Overview")

# Add text content
left = Inches(0.5)
top = Inches(1.5)
width = Inches(6)
height = Inches(5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

content = """Key Statistics:

• Total System Length: 2,340 miles (Minneapolis to Gulf)
• Lower Mississippi: Mile 0-233 (New Orleans to Baton Rouge)
• Facilities Tracked: 842+ across the system
• Louisiana Facilities: 52 major terminals

System Components:
• Export grain elevators
• Petroleum refineries & tank terminals
• Chemical plants & bulk terminals
• Intermodal facilities
• Mid-stream transfer points
• General cargo terminals"""

p = tf.paragraphs[0]
p.text = content
p.font.size = Pt(18)
p.space_after = Pt(12)

# Add stats box
left = Inches(7)
top = Inches(1.5)
width = Inches(5.5)
height = Inches(5)

textbox2 = slide.shapes.add_textbox(left, top, width, height)
tf2 = textbox2.text_frame
tf2.word_wrap = True

stats = f"""Facility Count by Type:

Tank Storage: {len(facilities[facilities['facility_type'] == 'Tank Storage'])}
Refineries: {len(facilities[facilities['facility_type'] == 'Refinery'])}
Bulk Terminals: {len(facilities[facilities['facility_type'] == 'Bulk Terminal'])}
Grain Elevators: {len(facilities[facilities['facility_type'] == 'Elevator'])}
General Cargo: {len(facilities[facilities['facility_type'] == 'General Cargo'])}
Mid-Stream: {len(facilities[facilities['facility_type'] == 'Mid-Stream'])}
Bulk Plants: {len(facilities[facilities['facility_type'] == 'Bulk Plant'])}"""

p2 = tf2.paragraphs[0]
p2.text = stats
p2.font.size = Pt(18)

# SLIDE 3: Facility Types - Detailed
slide = add_content_slide(prs, "Major Facility Types")

left = Inches(0.5)
top = Inches(1.5)
width = Inches(12)
height = Inches(5.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

facilities_detail = """
EXPORT GRAIN ELEVATORS
• Receive grain by barge from Midwest
• Load ocean-going vessels (Panamax, Handymax)
• Capacity: 4-10 million bushels per facility
• Loading rates: 50,000-80,000 bushels/hour
• Key operators: ADM, Bunge, Cargill, CHS, Zen-Noh

PETROLEUM REFINERIES & TANK TERMINALS
• Process crude oil into refined products
• Store crude oil, gasoline, diesel, jet fuel
• Capacity: 1-16 million barrels
• Deep-water vessel access for crude imports
• Key operators: Valero, Marathon, ExxonMobil, IMTT, Kinder Morgan

CHEMICAL TERMINALS & BULK PLANTS
• Store chemicals, vegetable oils, specialty products
• Heating/cooling/mixing capabilities
• Vapor control systems
• Key operators: Stolthaven, IMTT, International Raw Materials

BULK TERMINALS
• Coal, iron ore, steel products
• Aggregates, fertilizers
• Conveyor systems for loading/unloading
• Barge and ocean vessel capability"""

p = tf.paragraphs[0]
p.text = facilities_detail
p.font.size = Pt(16)
p.space_after = Pt(10)

# SLIDE 4: Vessel Types
slide = add_content_slide(prs, "Vessel Types Operating on the River")

left = Inches(0.5)
top = Inches(1.5)
width = Inches(12)
height = Inches(5.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

vessels = """
INLAND WATERWAY VESSELS

Towboats & Barge Tows
• Typical tow: 15-40 barges (jumbo barges: 195' x 35')
• Each barge: ~1,500 tons capacity
• Full tow capacity: 22,500-60,000 tons
• Primary cargo: grain, coal, petroleum products, chemicals, aggregates

OCEAN-GOING VESSELS

Panamax Vessels
• Length: ~700-750 feet
• Beam: 106 feet (Panama Canal max)
• Capacity: 60,000-80,000 DWT (deadweight tons)
• Draft: 39-45 feet (Mississippi River channel depth)
• Primary cargo: Export grain, bulk commodities

Handymax Vessels
• Length: ~550-600 feet
• Capacity: 40,000-60,000 DWT
• More flexible for smaller ports
• Common for grain, steel, forest products

Aframax Tankers (crude oil)
• Capacity: 80,000-120,000 DWT
• Deep-draft vessels for crude imports
• Access limited to lower river deepwater berths

Specialized Vessels
• Chemical tankers (stainless steel tanks)
• Product tankers (refined petroleum)
• Container ships (limited Mississippi operations)
• RoRo (roll-on/roll-off) for vehicles"""

p = tf.paragraphs[0]
p.text = vessels
p.font.size = Pt(15)

# SLIDE 5: Cargo Types
slide = add_content_slide(prs, "Major Cargo Categories")

left = Inches(0.5)
top = Inches(1.5)
width = Inches(5.5)
height = Inches(5.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

cargo1 = """AGRICULTURAL COMMODITIES

Export Grains
• Corn (largest by volume)
• Soybeans
• Wheat
• Rice
• Sorghum
• DDGS (distillers grains)

Annual throughput:
• Port of South Louisiana handles >50% of U.S. grain exports
• ~60 million metric tons/year

Major destinations:
• Asia Pacific (China, Japan, Korea)
• Latin America (Mexico, Colombia)
• Middle East & North Africa

PETROLEUM & REFINED PRODUCTS

Crude Oil
• Heavy Louisiana Sweet
• Light Louisiana Sweet
• Imported crude (Latin America, Middle East)

Refined Products
• Gasoline, diesel, jet fuel
• Heating oil, residual fuel
• Asphalt, petroleum coke
• Lubricants"""

p = tf.paragraphs[0]
p.text = cargo1
p.font.size = Pt(14)

# Right column
left = Inches(6.5)
textbox2 = slide.shapes.add_textbox(left, top, width, height)
tf2 = textbox2.text_frame
tf2.word_wrap = True

cargo2 = """CHEMICALS & SPECIALTY PRODUCTS

Bulk Chemicals
• Petrochemicals (benzene, toluene)
• Fertilizers (ammonia, urea)
• Caustic soda, acids
• Specialty chemicals

Vegetable Oils & Biodiesel
• Soybean oil
• Palm oil imports
• Renewable diesel feedstocks

BULK MATERIALS

Coal & Coke
• Steam coal for power generation
• Metallurgical coal for steel
• Petroleum coke

Metals & Ores
• Iron ore
• Steel products (coils, bars)
• Scrap metal
• Non-ferrous metals

Construction Materials
• Aggregates (sand, gravel, stone)
• Cement
• Steel for construction

Forest Products
• Lumber, logs
• Wood chips, pulp
• Paper products"""

p2 = tf2.paragraphs[0]
p2.text = cargo2
p2.font.size = Pt(14)

# SLIDE 6: Geographic Overview - Lower Mississippi
slide = add_content_slide(prs, "Lower Mississippi River - Key Zones")

left = Inches(0.5)
top = Inches(1.5)
width = Inches(12)
height = Inches(5.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

geography = """
ZONE 1: BELOW NEW ORLEANS (Mile 0-80)
• Chevron Empire (Mile 27) - Crude oil offshore delivery
• Phillips 66/Harvest Belle Chasse (Mile 63) - Crude terminal (former refinery)
• Stolthaven (Mile 80) - Chemical storage (2.1M bbl)
• Closest to Gulf of Mexico, deepwater access

ZONE 2: NEW ORLEANS METRO (Mile 80-103)
• Violet Dock, Domino Sugar, Port of NOLA terminals
• Harvey/Marrero/Gretna cluster:
  - IMTT Gretna (now BWC Terminals)
  - Kinder Morgan Harvey
  - Magellan/Buckeye terminals
• Container terminals, general cargo

ZONE 3: WESTWEGO TO DESTREHAN (Mile 103-155)
• Cargill Westwego (Mile 103) - Dual-berth grain elevator, 4.3M bu
• Zen-Noh Convent (Mile 140) - Grain elevator, 7.6M bu
• ADM Destrehan (Mile 147) - Largest grain elevator, 9.8M bu
• ADM Ama (Mile 121) - Grain elevator, 6M bu
• Grain export corridor

ZONE 4: RESERVE TO GRAMERCY (Mile 155-175)
• Marathon Garyville Refinery (Mile 160) - 578,000 bpd, largest in Louisiana
• Shell Norco (Mile 153-156) - Chemical complex
• Valero St. Charles (Mile 158) - Refinery
• IMTT Avondale (Mile 108) - Chemical storage
• Major refining & petrochemical hub

ZONE 5: BATON ROUGE AREA (Mile 175-233)
• Bunge Destrehan (Mile 158) - Grain elevator
• ExxonMobil Baton Rouge (Mile 228-230) - Refinery complex, 520,000 bpd
• IMTT Geismar (Mile 184) - Renewable fuels terminal
• Dow Plaquemine, Westlake, Shintech - Chemical plants
• Port of Greater Baton Rouge - Upriver terminal zone"""

p = tf.paragraphs[0]
p.text = geography
p.font.size = Pt(14)

# SLIDE 7: Facility Map (create and reference)
slide = add_content_slide(prs, "Facility Location Map - Interactive")

# Add text explaining the map
left = Inches(0.5)
top = Inches(1.5)
width = Inches(12)
height = Inches(1)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.text = "Facility distribution across the Mississippi River system. See interactive HTML maps in _html_web_files directory."
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.alignment = PP_ALIGN.CENTER

# Create the map
facility_map = create_facility_map(facilities)
map_file = save_map_as_image(facility_map, "_html_web_files/Facility_Overview_Map")

# Note: In production, you'd convert HTML to image and embed
# For now, we'll reference the HTML file
left = Inches(0.5)
top = Inches(3)
width = Inches(12)
height = Inches(0.8)

textbox2 = slide.shapes.add_textbox(left, top, width, height)
tf2 = textbox2.text_frame
tf2.text = f"Interactive map saved to: {map_file}"
p2 = tf2.paragraphs[0]
p2.font.size = Pt(14)
p2.font.italic = True
p2.alignment = PP_ALIGN.CENTER

# SLIDE 8: Commodity Flow Patterns
slide = add_content_slide(prs, "Commodity Flow Patterns")

left = Inches(0.5)
top = Inches(1.5)
width = Inches(12)
height = Inches(5.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

flows = """
EXPORT FLOWS (Outbound from Louisiana)

Grain Exports (Mississippi River → Gulf → World)
• Origin: Iowa, Illinois, Minnesota, Missouri grain belt
• Transport: 15-40 barge tows → Lower Miss elevators
• Destination: Asia Pacific (60%), Latin America (25%), Middle East/Africa (15%)
• Peak season: Post-harvest (Oct-Feb)
• Volume: 60+ million metric tons annually

Petroleum Product Exports
• Origin: Gulf Coast refineries
• Transport: Pipeline → tank terminal → ocean vessel
• Destinations: Latin America, Europe
• Products: Diesel, gasoline, jet fuel, petroleum coke

IMPORT FLOWS (Inbound to Louisiana)

Crude Oil Imports
• Origin: Mexico, Colombia, Saudi Arabia, Iraq
• Transport: Aframax tankers → deepwater terminals → pipeline
• Destination: Louisiana/Texas refineries
• Trend: Declining due to U.S. shale production

Specialty Chemicals & Oils
• Origin: Asia, Europe, South America
• Products: Vegetable oils, chemical feedstocks
• Storage: Tank terminals (Stolthaven, IMTT)

INTERNAL FLOWS (River-borne domestic commerce)

Coal (Upper Mississippi/Illinois → Louisiana power plants)
Steel Products (Great Lakes → Gulf construction markets)
Aggregates (Local quarries → construction sites)
Fertilizers (Gulf production → Midwest farms)
Chemicals (Louisiana plants → industrial users nationwide)"""

p = tf.paragraphs[0]
p.text = flows
p.font.size = Pt(14)

# SLIDE 9: Infrastructure Capacity Summary
slide = add_content_slide(prs, "System Capacity Summary")

left = Inches(0.5)
top = Inches(1.5)
width = Inches(5.5)
height = Inches(5.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

capacity1 = """GRAIN EXPORT CAPACITY

Louisiana Export Elevators (9 total)
• Total storage: ~64 million bushels
• Annual throughput: 60+ million metric tons
• Loading capacity: 50,000-80,000 bu/hr per facility

Top Facilities:
1. ADM Destrehan: 9.8M bu
2. Zen-Noh Convent: 7.6M bu
3. CHS Myrtle Grove: 6.5M bu
4. ADM Ama: 6M bu
5. Cargill Reserve: 5.5M bu

Vessel Loading:
• Panamax: 16-20 hours to load
• 2-4 vessels loaded per facility per week
• Peak season: Oct-Feb

PETROLEUM STORAGE

Tank Terminal Capacity
• IMTT St. Rose: 14.8M barrels (largest)
• Multiple operators: 40+ million bbl total
• Products: Crude, gasoline, diesel, chemicals"""

p = tf.paragraphs[0]
p.text = capacity1
p.font.size = Pt(13)

# Right column
left = Inches(6.5)
textbox2 = slide.shapes.add_textbox(left, top, width, height)
tf2 = textbox2.text_frame
tf2.word_wrap = True

capacity2 = """REFINERY CAPACITY

Louisiana Refineries (Top 3)
1. Marathon Garyville: 578,000 bpd
2. ExxonMobil Baton Rouge: 520,000 bpd
3. Valero Meraux: 340,000 bpd

Total Louisiana: ~2.9M bpd
(3rd largest refining state in U.S.)

WATERWAY CAPACITY

Navigation Channel
• Depth: 45 feet (SW Pass entrance)
• Width: 500-750 feet
• Controlling depth maintained by USACE

Vessel Traffic
• ~5,000 oceangoing vessels/year
• ~60,000 barge movements/year
• ~500,000 barges on inland system

Lock & Dam Capacity
• 15 locks on Upper Mississippi
• 8 locks on Ohio River
• Constrains upriver barge traffic

INTERMODAL CONNECTIONS

Rail Access
• Class 1: UP, BNSF, CN, KCS, NS, CSX
• Direct rail to most major terminals
• Unit trains for grain, coal, chemicals

Pipeline Network
• 50+ crude & product pipelines
• Connected to Cushing, OK hub
• Gulf Coast to Midwest corridors"""

p2 = tf2.paragraphs[0]
p2.text = capacity2
p2.font.size = Pt(13)

# SLIDE 10: Summary & Reference
slide = add_content_slide(prs, "Summary & Additional Resources")

left = Inches(0.5)
top = Inches(1.5)
width = Inches(12)
height = Inches(5.5)

textbox = slide.shapes.add_textbox(left, top, width, height)
tf = textbox.text_frame
tf.word_wrap = True

summary = """
KEY TAKEAWAYS

• Mississippi River system is the backbone of U.S. bulk commodity export infrastructure
• Louisiana Lower Mississippi (Mile 0-233) contains the world's largest concentration of export grain elevators
• System handles 500+ million tons of cargo annually across all commodity types
• Deep-draft navigation channel supports Panamax vessels up to Baton Rouge
• Multimodal connectivity: River + Rail + Pipeline creates integrated logistics network

FACILITY DISTRIBUTION

• Tank Storage & Refineries: Concentrated in Harvey/Marrero and Baton Rouge zones
• Grain Elevators: Westwego to Baton Rouge corridor (Mile 103-160)
• Chemical Terminals: Distributed throughout, major clusters at Norco/Garyville
• General Cargo: New Orleans metro area

ADDITIONAL RESOURCES

Interactive Maps Available:
• Lower_Mississippi_Industrial_Infrastructure_Map.html - All facilities with detailed popups
• National_Supply_Chain_Infrastructure_Map.html - Full U.S. waterway system (842 facilities)
• National_Market_Clusters_Map.html - Market accessibility analysis
• Commodity_Flows_Map.html - Commodity-specific flow visualization

Reports Available:
• Lower_Mississippi_River_Terminal_Facilities_Report.md - Detailed facility profiles
• Lower_Mississippi_River_Export_Grain_Elevators_Report.md - Grain elevator specifications
• Lower_Mississippi_River_Industrial_Infrastructure_Report.md - Comprehensive infrastructure overview

Data Files:
• master_facilities.csv - Complete facility registry with coordinates, commodities, operators
• national_supply_chain/ - National-scale geospatial datasets"""

p = tf.paragraphs[0]
p.text = summary
p.font.size = Pt(13)

# Save presentation
output_file = "_html_web_files/Mississippi_River_Facilities_Overview.pptx"
prs.save(output_file)

print(f"\n{'='*80}")
print(f"SUCCESS! Presentation created:")
print(f"  {output_file}")
print(f"\nSlides created:")
print(f"  1. Title Slide")
print(f"  2. System Overview")
print(f"  3. Major Facility Types")
print(f"  4. Vessel Types")
print(f"  5. Major Cargo Categories")
print(f"  6. Geographic Overview - Lower Mississippi")
print(f"  7. Facility Location Map")
print(f"  8. Commodity Flow Patterns")
print(f"  9. System Capacity Summary")
print(f"  10. Summary & Additional Resources")
print(f"\nInteractive map created:")
print(f"  _html_web_files/Facility_Overview_Map.html")
print(f"{'='*80}\n")

# Also create the standalone interactive map
facility_map.save("_html_web_files/Facility_Overview_Map.html")
print("Standalone interactive facility map saved!")
